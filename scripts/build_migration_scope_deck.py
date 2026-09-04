"""Generate the customer-facing "What can / cannot be migrated" slide deck.

The content is derived from what the tool actually does:

  * ``pages.BACKUP_ELEMENTS``      — everything the backup captures
  * ``RestorePage._run_restore``   — the branches that actually write to the
                                     destination
  * ``tests/test_restore_coverage._NOT_RESTORED`` — elements captured for
                                     reporting only and deliberately never
                                     written

Run (needs python-pptx, which the app itself does not depend on):
    pip install python-pptx
    python scripts/build_migration_scope_deck.py

Output:
    docs/S1-Migration-Scope-v<APP_VERSION>.pptx
"""
from __future__ import annotations

import datetime as _dt
import os
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
    _HAVE_PPTX = True
except ImportError:  # keep the content tables importable for the guard test
    _HAVE_PPTX = False
    Presentation = MSO_SHAPE = MSO_ANCHOR = PP_ALIGN = None

    def Emu(v):
        return int(v)

    def Inches(v):
        return int(v * 914400)

    def Pt(v):
        return int(v * 12700)

    def RGBColor(r, g, b):
        return (r, g, b)

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import APP_VERSION  # noqa: E402

# ── Brand ────────────────────────────────────────────────────────────────
PURPLE = RGBColor(0x6B, 0x0A, 0xEA)
INK = RGBColor(0x11, 0x12, 0x2B)
BODY = RGBColor(0x33, 0x35, 0x4D)
MUTED = RGBColor(0x76, 0x78, 0x92)
CARD_BG = RGBColor(0xF8, 0xF8, 0xFC)
CARD_LINE = RGBColor(0xE4, 0xE4, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.62)
GAP = Inches(0.4)
CARD_W = Emu(int((SLIDE_W - 2 * MARGIN - GAP) / 2))
CARD_Y = Inches(1.5)
CARD_H_MAX = Inches(5.1)
FOOTER_H = Inches(0.55)

# Hanging indent so wrapped bullet lines align with the text, not the dot.
BULLET_MAR_L = Inches(0.19)


# ── Primitives ───────────────────────────────────────────────────────────
def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def _footer(slide, note: str = ""):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "SentinelOne"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT

    if note:
        box = slide.shapes.add_textbox(
            SLIDE_W - Inches(6.0) - Inches(0.5), SLIDE_H - FOOTER_H,
            Inches(6.0), FOOTER_H)
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.RIGHT
        bp.text = note
        bp.font.size = Pt(10)
        bp.font.color.rgb = WHITE
        bp.font.name = FONT


def _title(slide, text: str, sub: str = ""):
    box = slide.shapes.add_textbox(MARGIN, Inches(0.42),
                                   SLIDE_W - 2 * MARGIN, Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = INK
    p.font.name = FONT
    if sub:
        sp = tf.add_paragraph()
        sp.text = sub
        sp.font.size = Pt(12.5)
        sp.font.color.rgb = MUTED
        sp.font.name = FONT
        sp.space_before = Pt(4)


def _card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.035
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


CARD_PAD = Inches(0.34)


def _hang(paragraph, mar_l=BULLET_MAR_L):
    """Give a paragraph a hanging indent (python-pptx exposes no API for it)."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(mar_l)))
    pPr.set("indent", str(-int(mar_l)))


def _wrapped_lines(text: str, pt: float, width_in: float) -> int:
    """Rough line count for Arial at `pt` inside `width_in` inches."""
    per_line = max(12, int(width_in / (0.55 * pt / 72.0)))
    words, lines, cur = text.split(), 1, 0
    for word in words:
        need = len(word) + (1 if cur else 0)
        if cur + need > per_line:
            lines += 1
            cur = len(word)
        else:
            cur += need
    return lines


def _card_height(card, body_pt: float) -> int:
    """Height needed for a card's content, so cards hug their bullets."""
    text_w = (CARD_W - 2 * CARD_PAD) / 914400.0
    total = 0.26 + 0.30  # top + bottom padding
    for i, (heading, bullets) in enumerate(card["sections"]):
        total += (12 / 72.0) if i else 0.0        # space_before
        total += (15 * 1.2 + 5) / 72.0            # heading + space_after
        for b in bullets:
            n = _wrapped_lines(b, body_pt, text_w - BULLET_MAR_L / 914400.0)
            total += (n * body_pt * 1.2 + 2.5) / 72.0
    if card.get("note"):
        n = _wrapped_lines(card["note"], 10, text_w)
        total += (10 + n * 10 * 1.2) / 72.0
    return min(int(Inches(total)), int(CARD_H_MAX))


def _fill_card(slide, x, y, w, h, sections, note="", body_pt=12.0):
    """sections = [(heading, [bullet, ...]), ...]"""
    _card(slide, x, y, w, h)
    box = slide.shapes.add_textbox(x + CARD_PAD, y + Inches(0.26),
                                   w - 2 * CARD_PAD, h - Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    for heading, bullets in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = heading
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.font.name = FONT
        p.space_before = Pt(0 if first else 12)
        p.space_after = Pt(5)
        first = False

        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = f"\u2022  {b}"
            bp.font.size = Pt(body_pt)
            bp.font.color.rgb = BODY
            bp.font.name = FONT
            bp.space_after = Pt(2.5)
            bp.line_spacing = 1.0
            _hang(bp)

    if note:
        np_ = tf.add_paragraph()
        np_.text = note
        np_.font.size = Pt(10)
        np_.font.italic = True
        np_.font.color.rgb = MUTED
        np_.font.name = FONT
        np_.space_before = Pt(10)


def _two_cards(slide, left, right, body_pt=12.0):
    lx, rx = MARGIN, MARGIN + CARD_W + GAP
    h = Emu(max(_card_height(left, body_pt), _card_height(right, body_pt)))
    _fill_card(slide, lx, CARD_Y, CARD_W, h,
               left["sections"], left.get("note", ""), body_pt)
    _fill_card(slide, rx, CARD_Y, CARD_W, h,
               right["sections"], right.get("note", ""), body_pt)


# ── Content (verified against pages.py) ──────────────────────────────────
MIGRATED_1_LEFT = {
    "sections": [
        ("Structure & Policies", [
            "Account, Site and Group structure (created on the destination)",
            "Group ranking \u2014 priority order preserved per site",
            "Policies at Global, Account, Site and Group scope",
            "Locations (firewall location-awareness)",
            "Saved Filters (Deep Visibility / SDL queries)",
            "Config Overrides (a.k.a. Policy Overrides)",
            "Agent Auto-Upgrade Policies",
            "Log Collection Rules (XDR ingestion)",
        ]),
    ],
    "note": "New accounts are provisioned with the destination tenant\u2019s "
            "primary licence bundle.",
}

MIGRATED_1_RIGHT = {
    "sections": [
        ("Security Controls", [
            "Exclusions \u2014 hash, path, file type, certificate, browser",
            "Unified Exclusions (v2.1), including tag-based exclusions",
            "Blocklist (SHA1 / SHA256)",
            "STAR Custom Detection Rules",
            "Threat Intelligence IOCs (account scope)",
        ]),
        ("Tags", [
            "Endpoint / Device Inventory tags (Tag Manager)",
            "Firewall Control rule tags",
            "Network Quarantine rule tags",
        ]),
    ],
}

MIGRATED_2_LEFT = {
    "sections": [
        ("Network & Device Protection", [
            "Firewall Control configuration",
            "Firewall Control rules",
            "Network Quarantine configuration",
            "Network Quarantine rules",
            "Device Control configuration",
            "Device Control rules",
        ]),
    ],
    "note": "Rule ordering is preserved during migration for Firewall Control "
            "and Device Control.",
}

MIGRATED_2_RIGHT = {
    "sections": [
        ("Settings & Integrations", [
            "SSO / SAML configuration",
            "SMTP relay settings \u2014 password excluded (see caveats)",
            "Syslog forwarding",
            "Active Directory integration",
            "Notification settings & recipients",
            "Scheduled reports",
        ]),
        ("Access & Identity", [
            "RBAC custom roles (account scope)",
            "Console users \u2014 locally-created users only",
            "Service users \u2014 re-created with their scope and roles; "
            "issue a new API token on the destination",
        ]),
    ],
}

NOT_LEFT = {
    "sections": [
        ("Listed in the report \u2014 re-create manually", [
            "API tokens \u2014 SentinelOne reveals a token once, at creation, "
            "so it is never in the backup",
            "Management gateways / proxies \u2014 environment-specific",
            "Singularity Marketplace integrations \u2014 each needs its own "
            "OAuth / credentials",
            "Remote Script Orchestration (RSO) script bodies \u2014 held in "
            "per-tenant storage",
            "Webhooks (Slack / Teams / generic HTTP) \u2014 no webhook "
            "endpoint exists in the v2.1 API",
            "SMTP password \u2014 write-only on the API",
        ]),
    ],
    "note": "These are captured in the backup and listed in the migration "
            "report so nothing is forgotten \u2014 they are simply not "
            "writable via the API.",
}

NOT_RIGHT = {
    "sections": [
        ("Not transferred", [
            "Expired / deleted accounts and sites (auto-skipped)",
            "Activity logs",
            "Historical incidents and threat data",
            "Deep Visibility event data \u2014 saved filters ARE migrated",
            "Agent data \u2014 agents stay on the source until moved "
            "separately",
            "Network Discovery data (Ranger)",
            "Cloud Native Security / Singularity Identity data",
            "MDR / Vigilance settings",
            "Anything at Global scope when no global-scope token is provided",
        ]),
    ],
}

CAVEAT_LEFT = {
    "sections": [
        ("Secrets & identity", [
            "SMTP: everything migrates except the password \u2014 re-enter it "
            "once on the destination",
            "SSO: the SAML certificate and metadata are bound to the source "
            "tenant URL and must be re-issued",
            "Console users: SSO / SCIM users auto-provision on first login; "
            "local users are created and receive a SentinelOne invitation "
            "email",
            "RBAC: only custom roles are created \u2014 predefined roles "
            "already exist on the destination",
            "Service users: the user, its scope and its roles are re-created, "
            "but its API token cannot be \u2014 issue a new one and update "
            "whatever integration used the old token",
        ]),
    ],
}

CAVEAT_RIGHT = {
    "sections": [
        ("Scope & licensing", [
            "Threat Intel, RBAC roles, console/service users, scripts and "
            "marketplace apps are account-scope only",
            "Global-scope elements require a global (MSSP) API token on both "
            "consoles",
            "The destination must license the matching modules (XDR, Ranger, "
            "Marketplace) or those elements return 404",
            "Re-runs are safe \u2014 items that already exist on the "
            "destination are skipped, not duplicated",
        ]),
    ],
    "note": "Every run produces a per-item migration report and a validation "
            "pass that compares source and destination.",
}

MIGRATED_CARDS = (MIGRATED_1_LEFT, MIGRATED_1_RIGHT,
                  MIGRATED_2_LEFT, MIGRATED_2_RIGHT)
NOT_MIGRATED_CARDS = (NOT_LEFT, NOT_RIGHT)
CAVEAT_CARDS = (CAVEAT_LEFT, CAVEAT_RIGHT)

# Elements whose restore branch only *lists* the items in the migration report
# — the destination write is impossible via the API (script bodies live in
# per-tenant storage; each marketplace app needs its own OAuth). They belong on
# the "will not be migrated" slide even though they have a restore branch.
INVENTORY_ONLY = {"marketplace_apps", "scripts"}

# Every key in ``pages.BACKUP_ELEMENTS`` must be represented on a slide, on the
# side that matches what the restore loop actually does. Maps the element to
# ``(group, distinctive phrase that must appear in that group's bullets)``.
# tests/test_deck_coverage.py fails if backup grows an element the deck doesn't
# mention, or if an element is advertised as migrated while the restore loop
# declares it not-restored.
ELEMENT_SLIDES = {
    "policy":                  ("migrated", "Policies at Global, Account"),
    "exclusions":              ("migrated", "Exclusions \u2014 hash, path"),
    "unified_exclusions":      ("migrated", "Unified Exclusions"),
    "blocklist":               ("migrated", "Blocklist (SHA1 / SHA256)"),
    "firewall_rules":          ("migrated", "Firewall Control rules"),
    "firewall_config":         ("migrated", "Firewall Control configuration"),
    "nq_config":               ("migrated", "Network Quarantine configuration"),
    "nq_rules":                ("migrated", "Network Quarantine rules"),
    "device_control_rules":    ("migrated", "Device Control rules"),
    "device_control_config":   ("migrated", "Device Control configuration"),
    "tags_firewall":           ("migrated", "Firewall Control rule tags"),
    "tags_network_quarantine": ("migrated", "Network Quarantine rule tags"),
    "tags_endpoint":           ("migrated", "Device Inventory tags"),
    "star_rules":              ("migrated", "STAR Custom Detection Rules"),
    "saved_filters":           ("migrated", "Saved Filters"),
    "threat_intel":            ("migrated", "Threat Intelligence IOCs"),
    "config_overrides":        ("migrated", "Config Overrides"),
    "log_collection_rules":    ("migrated", "Log Collection Rules"),
    "auto_upgrade_policies":   ("migrated", "Auto-Upgrade Policies"),
    "locations":               ("migrated", "Locations (firewall"),
    "settings_notifications":  ("migrated", "Notification settings"),
    "settings_sso":            ("migrated", "SSO / SAML configuration"),
    "settings_smtp":           ("migrated", "SMTP relay settings"),
    "settings_syslog":         ("migrated", "Syslog forwarding"),
    "settings_ad":             ("migrated", "Active Directory integration"),
    "webhooks":                ("not_migrated", "Webhooks (Slack"),
    "scheduled_reports":       ("migrated", "Scheduled reports"),
    "roles":                   ("migrated", "RBAC custom roles"),
    "console_users":           ("migrated", "Console users \u2014 locally"),
    "service_users":           ("migrated", "Service users \u2014 re-created"),
    "gateways":                ("not_migrated", "Management gateways"),
    "marketplace_apps":        ("not_migrated", "Singularity Marketplace"),
    "scripts":                 ("not_migrated", "Remote Script Orchestration"),
}


def bullets(cards) -> list[str]:
    """Every bullet (and card note) across a group of cards."""
    out = []
    for card in cards:
        for _heading, items in card["sections"]:
            out.extend(items)
        if card.get("note"):
            out.append(card["note"])
    return out


def build(out_path: str) -> str:
    if not _HAVE_PPTX:
        raise RuntimeError(
            "python-pptx is not installed in this interpreter. Install it "
            "and re-run:\n"
            "  pip install python-pptx\n"
            "  python scripts/build_migration_scope_deck.py")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    today = _dt.date.today().strftime("%B %d, %Y")
    stamp = f"S1 Command Center v{APP_VERSION}  \u00b7  {today}"

    # 1 — title
    s = _blank(prs)
    box = s.shapes.add_textbox(MARGIN, Inches(2.35),
                               SLIDE_W - 2 * MARGIN, Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Console Migration Scope"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = INK
    p.font.name = FONT
    p2 = tf.add_paragraph()
    p2.text = "What can \u2014 and cannot \u2014 be migrated between " \
              "SentinelOne consoles"
    p2.font.size = Pt(19)
    p2.font.color.rgb = PURPLE
    p2.font.name = FONT
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = stamp
    p3.font.size = Pt(12)
    p3.font.color.rgb = MUTED
    p3.font.name = FONT
    p3.space_before = Pt(16)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.15),
                              Inches(1.6), Pt(5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = PURPLE
    rule.line.fill.background()
    rule.shadow.inherit = False
    _footer(s)

    # 2 — migrated 1/2
    s = _blank(prs)
    _title(s, "What types of settings can be migrated?", "1 of 2")
    _two_cards(s, MIGRATED_1_LEFT, MIGRATED_1_RIGHT)
    _footer(s, stamp)

    # 3 — migrated 2/2
    s = _blank(prs)
    _title(s, "What types of settings can be migrated?", "2 of 2")
    _two_cards(s, MIGRATED_2_LEFT, MIGRATED_2_RIGHT)
    _footer(s, stamp)

    # 4 — not migrated
    s = _blank(prs)
    _title(s, "What will not be migrated?")
    _two_cards(s, NOT_LEFT, NOT_RIGHT, body_pt=11.5)
    _footer(s, stamp)

    # 5 — caveats
    s = _blank(prs)
    _title(s, "Before you migrate \u2014 key caveats")
    _two_cards(s, CAVEAT_LEFT, CAVEAT_RIGHT, body_pt=12.0)
    _footer(s, stamp)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    dest = os.path.join(root, "docs",
                        f"S1-Migration-Scope-v{APP_VERSION}.pptx")
    print(build(dest))
